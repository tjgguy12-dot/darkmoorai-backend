"""
Office Editor Module – Enhanced with Conversion, Translation, Resume & Invoice Builders
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime, timedelta
import uuid
import base64
import json
import asyncio
import csv
from io import BytesIO

# Document processing
from docx import Document
from docx.shared import Inches, Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# Excel processing
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from openpyxl.utils import get_column_letter

# PowerPoint processing
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN

# PDF processing
import PyPDF2
import pdfplumber
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# Image processing
from PIL import Image
import requests

# QR code
import qrcode

# HTTP for translation
import httpx

# Internal imports
from app.config import config
from app.utils.logger import logger

# ============================================================================
# Resume Builder Class
# ============================================================================

class ResumeBuilder:
    """Professional resume builder with photo and multiple templates"""
    
    def __init__(self, temp_dir):
        self.temp_dir = temp_dir
    
    def create_resume(self, data: dict, template: str = "modern") -> str:
        templates = {
            "modern": self._create_modern_resume,
            "classic": self._create_classic_resume,
            "creative": self._create_creative_resume,
            "executive": self._create_executive_resume,
            "tech": self._create_tech_resume
        }
        creator = templates.get(template, self._create_modern_resume)
        return creator(data)
    
    def _add_photo(self, doc, photo_data: dict, position: str = "right"):
        try:
            if photo_data.get("url"):
                response = requests.get(photo_data["url"], timeout=10)
                img = Image.open(BytesIO(response.content))
            elif photo_data.get("base64"):
                img_data = base64.b64decode(photo_data["base64"])
                img = Image.open(BytesIO(img_data))
            elif photo_data.get("path"):
                img = Image.open(photo_data["path"])
            else:
                return None
            img.thumbnail((150, 150))
            temp_img = self.temp_dir / f"photo_{uuid.uuid4().hex[:8]}.png"
            img.save(temp_img)
            if position == "right":
                table = doc.add_table(rows=1, cols=2)
                table.autofit = False
                table.columns[0].width = Inches(4.5)
                table.columns[1].width = Inches(1.5)
                name_cell = table.cell(0, 0)
                photo_cell = table.cell(0, 1)
                photo_cell.paragraphs[0].add_run().add_picture(str(temp_img), width=Inches(1.2))
                photo_cell.paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
                return table
            else:
                paragraph = doc.add_paragraph()
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                paragraph.add_run().add_picture(str(temp_img), width=Inches(1.5))
                return None
        except Exception as e:
            logger.warning(f"Photo addition failed: {e}")
            return None
    
    def _style_heading(self, heading, style: str):
        for run in heading.runs:
            run.font.bold = True
            if style == "modern":
                run.font.color.rgb = RGBColor(0, 112, 192)
    
    def _create_modern_resume(self, data: dict) -> str:
        doc = Document()
        sections = doc.sections
        for section in sections:
            section.top_margin = Cm(1.5)
            section.bottom_margin = Cm(1.5)
            section.left_margin = Cm(1.5)
            section.right_margin = Cm(1.5)
        
        if data.get("photo"):
            self._add_photo(doc, data["photo"], "right")
        
        name = doc.add_heading(data.get("name", "Your Name"), 0)
        name.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title = doc.add_paragraph(data.get("title", "Professional"))
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title.runs[0].italic = True
        
        contact = doc.add_paragraph()
        contact.alignment = WD_ALIGN_PARAGRAPH.CENTER
        contact_info = []
        if data.get("email"): contact_info.append(f"✉ {data['email']}")
        if data.get("phone"): contact_info.append(f"📞 {data['phone']}")
        if data.get("location"): contact_info.append(f"📍 {data['location']}")
        if data.get("linkedin"): contact_info.append(f"🔗 {data['linkedin']}")
        if data.get("github"): contact_info.append(f"💻 {data['github']}")
        contact.add_run(" | ".join(contact_info))
        doc.add_paragraph()
        
        table = doc.add_table(rows=1, cols=2)
        table.autofit = False
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(4.5)
        left_cell = table.cell(0, 0)
        right_cell = table.cell(0, 1)
        
        # LEFT COLUMN
        if data.get("skills"):
            heading = left_cell.add_heading("SKILLS", level=2)
            self._style_heading(heading, "modern")
            for skill in data["skills"]:
                p = left_cell.add_paragraph()
                p.add_run("▹ ").bold = True
                if isinstance(skill, dict):
                    p.add_run(f"{skill.get('name', '')}")
                    if skill.get("level"):
                        p.add_run(f" - {skill.get('level')}")
                else:
                    p.add_run(str(skill))
            left_cell.add_paragraph()
        
        if data.get("languages"):
            heading = left_cell.add_heading("LANGUAGES", level=2)
            self._style_heading(heading, "modern")
            for lang in data["languages"]:
                p = left_cell.add_paragraph()
                p.add_run("▹ ").bold = True
                if isinstance(lang, dict):
                    p.add_run(f"{lang.get('language', '')} - {lang.get('level', '')}")
                else:
                    p.add_run(str(lang))
            left_cell.add_paragraph()
        
        if data.get("certifications"):
            heading = left_cell.add_heading("CERTIFICATIONS", level=2)
            self._style_heading(heading, "modern")
            for cert in data["certifications"]:
                p = left_cell.add_paragraph()
                p.add_run("▹ ").bold = True
                if isinstance(cert, dict):
                    p.add_run(f"{cert.get('name', '')}")
                    if cert.get("date"):
                        p.add_run(f" ({cert.get('date')})")
                else:
                    p.add_run(str(cert))
            left_cell.add_paragraph()
        
        # RIGHT COLUMN
        if data.get("summary"):
            heading = right_cell.add_heading("PROFESSIONAL SUMMARY", level=2)
            self._style_heading(heading, "modern")
            right_cell.add_paragraph(data["summary"])
            right_cell.add_paragraph()
        
        if data.get("experience"):
            heading = right_cell.add_heading("WORK EXPERIENCE", level=2)
            self._style_heading(heading, "modern")
            for exp in data["experience"]:
                p = right_cell.add_paragraph()
                p.add_run(exp.get("title", "")).bold = True
                p.add_run(f" - {exp.get('company', '')}")
                p = right_cell.add_paragraph()
                p.add_run(exp.get("date", "")).italic = True
                right_cell.add_paragraph(exp.get("description", ""))
                if exp.get("achievements"):
                    for ach in exp["achievements"]:
                        right_cell.add_paragraph(f"• {ach}", style='List Bullet')
            right_cell.add_paragraph()
        
        if data.get("education"):
            heading = right_cell.add_heading("EDUCATION", level=2)
            self._style_heading(heading, "modern")
            for edu in data["education"]:
                p = right_cell.add_paragraph()
                p.add_run(edu.get("degree", "")).bold = True
                p.add_run(f" - {edu.get('school', '')}")
                if edu.get("date"):
                    right_cell.add_paragraph(edu.get("date")).italic = True
                if edu.get("gpa"):
                    right_cell.add_paragraph(f"GPA: {edu.get('gpa')}")
                if edu.get("courses"):
                    right_cell.add_paragraph(f"Relevant Courses: {', '.join(edu['courses'])}")
            right_cell.add_paragraph()
        
        if data.get("projects"):
            heading = right_cell.add_heading("PROJECTS", level=2)
            self._style_heading(heading, "modern")
            for proj in data["projects"]:
                p = right_cell.add_paragraph()
                p.add_run(proj.get("name", "")).bold = True
                if proj.get("technologies"):
                    p.add_run(f" - {proj.get('technologies')}")
                right_cell.add_paragraph(proj.get("description", ""))
                if proj.get("link"):
                    right_cell.add_paragraph(f"Link: {proj.get('link')}")
        
        filename = f"resume_{data.get('name', 'resume').replace(' ', '_')}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def _create_classic_resume(self, data: dict) -> str:
        doc = Document()
        name = doc.add_heading(data.get("name", "Your Name"), 0)
        name.alignment = WD_ALIGN_PARAGRAPH.CENTER
        contact = doc.add_paragraph()
        contact.alignment = WD_ALIGN_PARAGRAPH.CENTER
        contact_info = []
        if data.get("email"): contact_info.append(data['email'])
        if data.get("phone"): contact_info.append(data['phone'])
        if data.get("location"): contact_info.append(data['location'])
        contact.add_run(" | ".join(contact_info))
        doc.add_paragraph("_" * 60)
        if data.get("summary"):
            doc.add_heading("SUMMARY", level=1)
            doc.add_paragraph(data["summary"])
        if data.get("experience"):
            doc.add_heading("EXPERIENCE", level=1)
            for exp in data["experience"]:
                p = doc.add_paragraph()
                p.add_run(exp.get("title", "")).bold = True
                doc.add_paragraph(f"{exp.get('company', '')} | {exp.get('date', '')}")
                doc.add_paragraph(exp.get("description", ""))
        if data.get("education"):
            doc.add_heading("EDUCATION", level=1)
            for edu in data["education"]:
                p = doc.add_paragraph()
                p.add_run(edu.get("degree", "")).bold = True
                doc.add_paragraph(edu.get("school", ""))
                if edu.get("date"):
                    doc.add_paragraph(edu.get("date"))
        if data.get("skills"):
            doc.add_heading("SKILLS", level=1)
            skills_text = ", ".join([s if isinstance(s, str) else s.get('name', '') for s in data["skills"]])
            doc.add_paragraph(skills_text)
        filename = f"resume_classic_{data.get('name', 'resume').replace(' ', '_')}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def _create_creative_resume(self, data: dict) -> str:
        doc = Document()
        doc.add_paragraph("_" * 80)
        name = doc.add_heading(data.get("name", "Your Name"), 0)
        name.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in name.runs:
            run.font.color.rgb = RGBColor(41, 128, 185)
        title = doc.add_paragraph(data.get("title", "Creative Professional"))
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title.runs[0].font.color.rgb = RGBColor(52, 73, 94)
        if data.get("portfolio"):
            doc.add_paragraph(f"Portfolio: {data['portfolio']}").alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()
        if data.get("skills"):
            doc.add_heading("Core Competencies", level=1)
            for skill in data["skills"]:
                skill_name = skill if isinstance(skill, str) else skill.get('name', '')
                level_val = 70
                if isinstance(skill, dict) and skill.get("level"):
                    level_map = {"Beginner": 30, "Intermediate": 60, "Advanced": 85, "Expert": 100}
                    level_val = level_map.get(skill["level"], 70)
                p = doc.add_paragraph()
                p.add_run(f"{skill_name}: ").bold = True
                p.add_run("█" * (level_val // 10) + "░" * (10 - (level_val // 10)))
        doc.add_paragraph()
        if data.get("experience"):
            doc.add_heading("Selected Achievements", level=1)
            for exp in data["experience"]:
                doc.add_heading(exp.get("title", ""), level=2)
                doc.add_paragraph(f"{exp.get('company', '')} | {exp.get('date', '')}")
                if exp.get("achievements"):
                    for ach in exp["achievements"]:
                        doc.add_paragraph(f"✨ {ach}", style='List Bullet')
        filename = f"resume_creative_{data.get('name', 'resume').replace(' ', '_')}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def _create_executive_resume(self, data: dict) -> str:
        doc = Document()
        name = doc.add_heading(data.get("name", "Your Name"), 0)
        title = doc.add_paragraph(data.get("title", "Executive Leader"))
        title.runs[0].bold = True
        contact = doc.add_paragraph()
        contact_info = []
        if data.get("email"): contact_info.append(data['email'])
        if data.get("phone"): contact_info.append(data['phone'])
        if data.get("linkedin"): contact_info.append(data['linkedin'])
        contact.add_run(" | ".join(contact_info))
        doc.add_paragraph()
        if data.get("summary"):
            doc.add_heading("EXECUTIVE SUMMARY", level=1)
            doc.add_paragraph(data["summary"])
        if data.get("experience"):
            doc.add_heading("LEADERSHIP EXPERIENCE", level=1)
            for exp in data["experience"]:
                p = doc.add_paragraph()
                p.add_run(exp.get("title", "")).bold = True
                p.add_run(f" | {exp.get('company', '')}")
                doc.add_paragraph(exp.get("date", "")).italic = True
                doc.add_paragraph(exp.get("description", ""))
                if exp.get("metrics"):
                    for metric in exp["metrics"]:
                        doc.add_paragraph(f"✓ {metric}", style='List Bullet')
        if data.get("skills"):
            doc.add_heading("CORE COMPETENCIES", level=1)
            skills_list = [s if isinstance(s, str) else s.get('name', '') for s in data["skills"]]
            for i in range(0, len(skills_list), 3):
                doc.add_paragraph(" | ".join(skills_list[i:i+3]))
        filename = f"resume_executive_{data.get('name', 'resume').replace(' ', '_')}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def _create_tech_resume(self, data: dict) -> str:
        doc = Document()
        name = doc.add_heading(data.get("name", "Your Name"), 0)
        title = doc.add_paragraph(data.get("title", "Software Engineer"))
        title.runs[0].bold = True
        contact = doc.add_paragraph()
        contact_info = []
        if data.get("email"): contact_info.append(f"✉ {data['email']}")
        if data.get("phone"): contact_info.append(f"📱 {data['phone']}")
        if data.get("github"): contact_info.append(f"🐙 {data['github']}")
        if data.get("linkedin"): contact_info.append(f"🔗 {data['linkedin']}")
        contact.add_run("  |  ".join(contact_info))
        doc.add_paragraph()
        if data.get("skills"):
            doc.add_heading("TECHNICAL SKILLS", level=1)
            skill_categories = {"Languages": [], "Frameworks": [], "Tools": [], "Databases": []}
            for skill in data["skills"]:
                if isinstance(skill, dict):
                    cat = skill.get("category", "Languages")
                    if cat in skill_categories:
                        skill_categories[cat].append(skill.get("name", ""))
                else:
                    skill_categories["Languages"].append(skill)
            for cat, skills in skill_categories.items():
                if skills:
                    doc.add_heading(cat, level=2)
                    doc.add_paragraph(", ".join(skills))
        if data.get("experience"):
            doc.add_heading("EXPERIENCE", level=1)
            for exp in data["experience"]:
                p = doc.add_paragraph()
                p.add_run(exp.get("title", "")).bold = True
                p.add_run(f" @ {exp.get('company', '')}")
                doc.add_paragraph(exp.get("date", "")).italic = True
                doc.add_paragraph(exp.get("description", ""))
                if exp.get("technologies"):
                    doc.add_paragraph(f"Tech Stack: {exp.get('technologies')}")
        if data.get("projects"):
            doc.add_heading("PROJECTS", level=1)
            for proj in data["projects"]:
                p = doc.add_paragraph()
                p.add_run(proj.get("name", "")).bold = True
                doc.add_paragraph(proj.get("description", ""))
                if proj.get("link"):
                    doc.add_paragraph(f"🔗 {proj.get('link')}")
        filename = f"resume_tech_{data.get('name', 'resume').replace(' ', '_')}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)


# ============================================================================
# Invoice Builder Class
# ============================================================================

class InvoiceBuilder:
    """Professional invoice builder with logo and multiple templates"""
    
    def __init__(self, temp_dir):
        self.temp_dir = temp_dir
    
    def create_invoice(self, data: dict, template: str = "professional") -> str:
        templates = {
            "professional": self._create_professional_invoice,
            "modern": self._create_modern_invoice,
            "minimal": self._create_minimal_invoice,
            "corporate": self._create_corporate_invoice,
            "simple": self._create_simple_invoice
        }
        creator = templates.get(template, self._create_professional_invoice)
        return creator(data)
    
    def _add_logo(self, doc, logo_data: dict, position: str = "left"):
        try:
            if logo_data.get("url"):
                response = requests.get(logo_data["url"], timeout=10)
                img = Image.open(BytesIO(response.content))
            elif logo_data.get("base64"):
                img_data = base64.b64decode(logo_data["base64"])
                img = Image.open(BytesIO(img_data))
            elif logo_data.get("path"):
                img = Image.open(logo_data["path"])
            else:
                return None
            img.thumbnail((120, 120))
            temp_img = self.temp_dir / f"logo_{uuid.uuid4().hex[:8]}.png"
            img.save(temp_img)
            paragraph = doc.add_paragraph()
            if position == "center":
                paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            paragraph.add_run().add_picture(str(temp_img), width=Inches(1.2))
            return paragraph
        except Exception as e:
            logger.warning(f"Logo addition failed: {e}")
            return None
    
    def _add_qr_code(self, doc, payment_data: dict):
        try:
            qr_data = f"Pay to: {payment_data.get('bank_name', '')}\n"
            qr_data += f"Account: {payment_data.get('account_number', '')}\n"
            qr_data += f"Amount: {payment_data.get('total', '')}\n"
            qr_data += f"Reference: INV-{payment_data.get('invoice_number', '')}"
            qr = qrcode.make(qr_data)
            temp_qr = self.temp_dir / f"qr_{uuid.uuid4().hex[:8]}.png"
            qr.save(temp_qr)
            paragraph = doc.add_paragraph()
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            paragraph.add_run().add_picture(str(temp_qr), width=Inches(1.5))
            doc.add_paragraph("Scan to pay")
            return paragraph
        except Exception as e:
            logger.warning(f"QR generation failed: {e}")
            return None
    
    def _create_professional_invoice(self, data: dict) -> str:
        doc = Document()
        if data.get("company_logo"):
            self._add_logo(doc, data["company_logo"], "left")
        company = doc.add_paragraph()
        company.alignment = WD_ALIGN_PARAGRAPH.CENTER
        if data.get("company_name"):
            company.add_run(data["company_name"]).bold = True
            company.add_run("\n")
        if data.get("company_address"):
            company.add_run(data["company_address"])
            company.add_run("\n")
        if data.get("company_email") or data.get("company_phone"):
            contact = []
            if data.get("company_email"): contact.append(data["company_email"])
            if data.get("company_phone"): contact.append(data["company_phone"])
            if data.get("company_website"): contact.append(data["company_website"])
            company.add_run(" | ".join(contact))
        doc.add_paragraph()
        title = doc.add_heading("INVOICE", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title.runs:
            run.font.color.rgb = RGBColor(0, 112, 192)
        table = doc.add_table(rows=2, cols=2)
        table.style = 'Table Grid'
        bill_to = table.cell(0, 0)
        bill_to.paragraphs[0].add_run("BILL TO:").bold = True
        if data.get("client_name"): bill_to.add_paragraph(data["client_name"])
        if data.get("client_email"): bill_to.add_paragraph(data["client_email"])
        if data.get("client_address"): bill_to.add_paragraph(data["client_address"])
        if data.get("client_phone"): bill_to.add_paragraph(data["client_phone"])
        invoice_info = table.cell(0, 1)
        invoice_info.paragraphs[0].add_run("INVOICE INFO:").bold = True
        invoice_number = data.get("invoice_number", f"INV-{datetime.now().strftime('%Y%m%d')}-001")
        invoice_info.add_paragraph(f"Invoice #: {invoice_number}")
        invoice_info.add_paragraph(f"Date: {data.get('invoice_date', datetime.now().strftime('%Y-%m-%d'))}")
        invoice_info.add_paragraph(f"Due Date: {data.get('due_date', (datetime.now() + timedelta(days=30)).strftime('%Y-%m-%d'))}")
        if data.get("po_number"): invoice_info.add_paragraph(f"PO Number: {data['po_number']}")
        payment_terms = table.cell(1, 0)
        payment_terms.paragraphs[0].add_run("PAYMENT TERMS:").bold = True
        payment_terms.add_paragraph(data.get("payment_terms", "Net 30 days"))
        bank_info = table.cell(1, 1)
        bank_info.paragraphs[0].add_run("BANK DETAILS:").bold = True
        if data.get("bank_name"): bank_info.add_paragraph(f"Bank: {data['bank_name']}")
        if data.get("account_name"): bank_info.add_paragraph(f"Account Name: {data['account_name']}")
        if data.get("account_number"): bank_info.add_paragraph(f"Account Number: {data['account_number']}")
        if data.get("routing_number"): bank_info.add_paragraph(f"Routing #: {data['routing_number']}")
        doc.add_paragraph()
        doc.add_heading("Items", level=2)
        items_table = doc.add_table(rows=1, cols=5)
        items_table.style = 'Table Grid'
        headers = ["#", "Description", "Quantity", "Unit Price", "Total"]
        for i, h in enumerate(headers):
            items_table.rows[0].cells[i].text = h
            items_table.rows[0].cells[i].paragraphs[0].runs[0].bold = True
        subtotal = 0
        for idx, item in enumerate(data.get("items", []), 1):
            row = items_table.add_row().cells
            row[0].text = str(idx)
            row[1].text = item.get("description", "")
            row[2].text = str(item.get("quantity", 1))
            row[3].text = f"${item.get('unit_price', 0):,.2f}"
            item_total = item.get("quantity", 1) * item.get("unit_price", 0)
            row[4].text = f"${item_total:,.2f}"
            subtotal += item_total
        doc.add_paragraph()
        totals_table = doc.add_table(rows=5, cols=2)
        totals_table.autofit = False
        totals_table.columns[0].width = Inches(5)
        totals_table.columns[1].width = Inches(2)
        tax_rate = data.get("tax_rate", 0)
        tax_amount = subtotal * (tax_rate / 100)
        discount = data.get("discount", 0)
        shipping = data.get("shipping", 0)
        total = subtotal + tax_amount - discount + shipping
        totals_table.cell(0, 0).text = "Subtotal:"
        totals_table.cell(0, 1).text = f"${subtotal:,.2f}"
        totals_table.cell(1, 0).text = f"Tax ({tax_rate}%):"
        totals_table.cell(1, 1).text = f"${tax_amount:,.2f}"
        totals_table.cell(2, 0).text = "Discount:"
        totals_table.cell(2, 1).text = f"-${discount:,.2f}"
        totals_table.cell(3, 0).text = "Shipping:"
        totals_table.cell(3, 1).text = f"${shipping:,.2f}"
        totals_table.cell(4, 0).text = "TOTAL:"
        totals_table.cell(4, 1).text = f"${total:,.2f}"
        totals_table.cell(4, 0).paragraphs[0].runs[0].bold = True
        totals_table.cell(4, 1).paragraphs[0].runs[0].bold = True
        if data.get("enable_qr"):
            payment_data = {
                "bank_name": data.get("bank_name", ""),
                "account_number": data.get("account_number", ""),
                "total": total,
                "invoice_number": invoice_number
            }
            self._add_qr_code(doc, payment_data)
        if data.get("notes"):
            doc.add_paragraph()
            doc.add_heading("Notes", level=2)
            doc.add_paragraph(data["notes"])
        footer = doc.add_paragraph()
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer.add_run("Thank you for your business!").italic = True
        filename = f"invoice_{invoice_number}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def _create_modern_invoice(self, data: dict) -> str:
        doc = Document()
        doc.add_paragraph("_" * 80)
        if data.get("company_logo"):
            self._add_logo(doc, data["company_logo"], "left")
        if data.get("company_name"):
            name = doc.add_heading(data["company_name"], level=1)
            for run in name.runs:
                run.font.color.rgb = RGBColor(41, 128, 185)
        title = doc.add_heading("INVOICE", level=2)
        for run in title.runs:
            run.font.color.rgb = RGBColor(231, 76, 60)
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).paragraphs[0].add_run("Client:").bold = True
        if data.get("client_name"): table.cell(0, 0).add_paragraph(data["client_name"])
        if data.get("client_email"): table.cell(0, 0).add_paragraph(data["client_email"])
        table.cell(0, 1).paragraphs[0].add_run("Invoice Details:").bold = True
        invoice_number = data.get("invoice_number", f"INV-{datetime.now().strftime('%Y%m%d')}-001")
        table.cell(0, 1).add_paragraph(f"#: {invoice_number}")
        table.cell(0, 1).add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d')}")
        items_table = doc.add_table(rows=1, cols=4)
        items_table.style = 'Table Grid'
        headers = ["Item", "Qty", "Price", "Total"]
        for i, h in enumerate(headers):
            items_table.rows[0].cells[i].text = h
            items_table.rows[0].cells[i].paragraphs[0].runs[0].bold = True
        for item in data.get("items", []):
            row = items_table.add_row().cells
            row[0].text = item.get("description", "")
            row[1].text = str(item.get("quantity", 1))
            row[2].text = f"${item.get('unit_price', 0):,.2f}"
            total = item.get("quantity", 1) * item.get("unit_price", 0)
            row[3].text = f"${total:,.2f}"
        filename = f"invoice_modern_{invoice_number}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def _create_minimal_invoice(self, data: dict) -> str:
        doc = Document()
        if data.get("company_name"):
            doc.add_heading(data["company_name"], level=0)
        else:
            doc.add_heading("INVOICE", level=0)
        doc.add_paragraph()
        table = doc.add_table(rows=1, cols=2)
        table.cell(0, 0).add_paragraph(f"Bill to:\n{data.get('client_name', '')}\n{data.get('client_email', '')}")
        invoice_number = data.get("invoice_number", f"INV-{datetime.now().strftime('%Y%m%d')}-001")
        table.cell(0, 1).add_paragraph(f"Invoice #: {invoice_number}\nDate: {datetime.now().strftime('%Y-%m-%d')}")
        items_table = doc.add_table(rows=1, cols=3)
        headers = ["Description", "Qty", "Amount"]
        for i, h in enumerate(headers):
            items_table.rows[0].cells[i].text = h
        for item in data.get("items", []):
            row = items_table.add_row().cells
            row[0].text = item.get("description", "")
            row[1].text = str(item.get("quantity", 1))
            total = item.get("quantity", 1) * item.get("unit_price", 0)
            row[2].text = f"${total:,.2f}"
        total = sum(item.get("quantity", 1) * item.get("unit_price", 0) for item in data.get("items", []))
        doc.add_paragraph(f"Total: ${total:,.2f}")
        filename = f"invoice_minimal_{invoice_number}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def _create_corporate_invoice(self, data: dict) -> str:
        doc = Document()
        header_table = doc.add_table(rows=1, cols=2)
        header_table.style = 'Table Grid'
        left_cell = header_table.cell(0, 0)
        if data.get("company_logo"):
            self._add_logo(doc, data["company_logo"], "left")
        if data.get("company_name"):
            left_cell.add_paragraph(data["company_name"]).bold = True
        if data.get("company_address"):
            left_cell.add_paragraph(data["company_address"])
        right_cell = header_table.cell(0, 1)
        right_cell.paragraphs[0].add_run("INVOICE").bold = True
        invoice_number = data.get("invoice_number", f"INV-{datetime.now().strftime('%Y%m%d')}-001")
        right_cell.add_paragraph(f"Number: {invoice_number}")
        right_cell.add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d')}")
        doc.add_paragraph()
        doc.add_heading("CLIENT INFORMATION", level=2)
        doc.add_paragraph(f"Name: {data.get('client_name', '')}")
        doc.add_paragraph(f"Email: {data.get('client_email', '')}")
        if data.get("client_phone"):
            doc.add_paragraph(f"Phone: {data['client_phone']}")
        doc.add_heading("INVOICE DETAILS", level=2)
        items_table = doc.add_table(rows=1, cols=5)
        items_table.style = 'Table Grid'
        headers = ["#", "Item Description", "Quantity", "Unit Price", "Line Total"]
        for i, h in enumerate(headers):
            items_table.rows[0].cells[i].text = h
        for idx, item in enumerate(data.get("items", []), 1):
            row = items_table.add_row().cells
            row[0].text = str(idx)
            row[1].text = item.get("description", "")
            row[2].text = str(item.get("quantity", 1))
            row[3].text = f"${item.get('unit_price', 0):,.2f}"
            total = item.get("quantity", 1) * item.get("unit_price", 0)
            row[4].text = f"${total:,.2f}"
        filename = f"invoice_corporate_{invoice_number}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def _create_simple_invoice(self, data: dict) -> str:
        doc = Document()
        doc.add_heading("INVOICE", level=0)
        invoice_number = data.get("invoice_number", f"INV-{datetime.now().strftime('%Y%m%d')}-001")
        doc.add_paragraph(f"Invoice #: {invoice_number}")
        doc.add_paragraph(f"Date: {datetime.now().strftime('%Y-%m-%d')}")
        doc.add_paragraph()
        doc.add_paragraph(f"Client: {data.get('client_name', '')}")
        doc.add_paragraph()
        for item in data.get("items", []):
            total = item.get("quantity", 1) * item.get("unit_price", 0)
            doc.add_paragraph(f"{item.get('description', '')} - {item.get('quantity', 1)} x ${item.get('unit_price', 0):,.2f} = ${total:,.2f}")
        total = sum(item.get("quantity", 1) * item.get("unit_price", 0) for item in data.get("items", []))
        doc.add_paragraph()
        doc.add_paragraph(f"TOTAL: ${total:,.2f}").bold = True
        filename = f"invoice_simple_{invoice_number}.docx"
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)


# ============================================================================
# Main OfficeEditor Class
# ============================================================================

class OfficeEditor:
    """Office document editor suite with conversion, translation, resume & invoice builders"""
    
    def __init__(self):
        self.temp_dir = Path("./data/office_temp")
        self.temp_dir.mkdir(parents=True, exist_ok=True)
        self.deepseek_api_key = config.DEEPSEEK_API_KEY
        self.deepseek_base_url = config.DEEPSEEK_BASE_URL
        self.resume_builder = ResumeBuilder(self.temp_dir)
        self.invoice_builder = InvoiceBuilder(self.temp_dir)
        logger.info("Office Editor initialized")
    
    # ============================================================================
    # Resume & Invoice Builders
    # ============================================================================
    
    def create_resume_with_template(self, data: dict, template: str = "modern") -> str:
        return self.resume_builder.create_resume(data, template)
    
    def create_invoice_with_template(self, data: dict, template: str = "professional") -> str:
        return self.invoice_builder.create_invoice(data, template)
    
    # ============================================================================
    # PDF Conversion Methods
    # ============================================================================
    
    def convert_pdf_to_word(self, pdf_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(pdf_path).stem}.docx"
        doc = Document()
        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    doc.add_heading(f"Page {page_num}", level=2)
                    doc.add_paragraph(text)
                    doc.add_page_break()
        doc.save(output_path)
        return str(output_path)
    
    def convert_pdf_to_excel(self, pdf_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(pdf_path).stem}.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "PDF Content"
        with pdfplumber.open(pdf_path) as pdf:
            row_num = 1
            for page_num, page in enumerate(pdf.pages, 1):
                ws.cell(row=row_num, column=1, value=f"=== Page {page_num} ===")
                row_num += 1
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        for col_idx, cell in enumerate(row, 1):
                            ws.cell(row=row_num, column=col_idx, value=cell if cell else "")
                        row_num += 1
                    row_num += 1
                text = page.extract_text()
                if text:
                    for line in text.split('\n'):
                        ws.cell(row=row_num, column=1, value=line)
                        row_num += 1
                row_num += 1
        wb.save(output_path)
        return str(output_path)
    
    def convert_pdf_to_powerpoint(self, pdf_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(pdf_path).stem}.pptx"
        prs = Presentation()
        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(reader.pages, 1):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                title.text = f"Page {page_num}"
                content = slide.placeholders[1]
                text_frame = content.text_frame
                text_frame.text = page.extract_text()[:500]
        prs.save(output_path)
        return str(output_path)
    
    # ============================================================================
    # Word Conversion Methods
    # ============================================================================
    
    def convert_word_to_pdf(self, docx_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(docx_path).stem}.pdf"
        doc = Document(docx_path)
        text_content = [para.text for para in doc.paragraphs if para.text.strip()]
        pdf_doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        for para in text_content:
            story.append(Paragraph(para, styles['Normal']))
            story.append(Spacer(1, 12))
        pdf_doc.build(story)
        return str(output_path)
    
    def convert_word_to_excel(self, docx_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(docx_path).stem}.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Document Content"
        doc = Document(docx_path)
        row_num = 1
        for para in doc.paragraphs:
            if para.text.strip():
                ws.cell(row=row_num, column=1, value=para.text)
                row_num += 1
        ws.column_dimensions['A'].width = 50
        wb.save(output_path)
        return str(output_path)
    
    def convert_word_to_csv(self, docx_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(docx_path).stem}.csv"
        doc = Document(docx_path)
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            for para in doc.paragraphs:
                if para.text.strip():
                    writer.writerow([para.text])
        return str(output_path)
    
    # ============================================================================
    # Excel Conversion Methods
    # ============================================================================
    
    def convert_excel_to_pdf(self, excel_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(excel_path).stem}.pdf"
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        ws = wb.active
        data = [[str(cell) if cell else "" for cell in row] for row in ws.iter_rows(values_only=True)]
        pdf_doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        if data:
            table = Table(data)
            table.setStyle(TableStyle([
                ('GRID', (0,0), (-1,-1), 0.5, colors.grey),
                ('FONTNAME', (0,0), (-1,-1), 'Helvetica'),
                ('FONTSIZE', (0,0), (-1,-1), 8),
            ]))
            story.append(table)
        pdf_doc.build(story)
        return str(output_path)
    
    def convert_excel_to_csv(self, excel_path: str, sheet_name: str = None) -> str:
        output_path = self.temp_dir / f"{Path(excel_path).stem}.csv"
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        ws = wb[sheet_name] if sheet_name else wb.active
        with open(output_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow(row)
        return str(output_path)
    
    def convert_excel_to_word(self, excel_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(excel_path).stem}.docx"
        doc = Document()
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            doc.add_heading(f"Sheet: {sheet_name}", level=1)
            data = [[str(cell) if cell else "" for cell in row] for row in ws.iter_rows(values_only=True)]
            if data:
                table = doc.add_table(rows=len(data), cols=len(data[0]))
                table.style = 'Table Grid'
                for i, row in enumerate(data):
                    for j, cell in enumerate(row):
                        table.cell(i, j).text = cell
            doc.add_page_break()
        doc.save(output_path)
        return str(output_path)
    
    # ============================================================================
    # PowerPoint Conversion Methods
    # ============================================================================
    
    def convert_powerpoint_to_pdf(self, pptx_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(pptx_path).stem}.pdf"
        prs = Presentation(pptx_path)
        content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            content.append(f"Slide {slide_num}")
            if slide.shapes.title:
                content.append(f"Title: {slide.shapes.title.text}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    content.append(shape.text)
            content.append("")
        pdf_doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        for line in content:
            story.append(Paragraph(line, styles['Normal']))
            story.append(Spacer(1, 6))
        pdf_doc.build(story)
        return str(output_path)
    
    def convert_powerpoint_to_word(self, pptx_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(pptx_path).stem}.docx"
        doc = Document()
        prs = Presentation(pptx_path)
        for slide_num, slide in enumerate(prs.slides, 1):
            doc.add_heading(f"Slide {slide_num}", level=1)
            if slide.shapes.title:
                doc.add_heading(slide.shapes.title.text, level=2)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    doc.add_paragraph(shape.text)
            doc.add_page_break()
        doc.save(output_path)
        return str(output_path)
    
    # ============================================================================
    # Smart Converter & Helpers
    # ============================================================================
    
    def convert_document(self, input_path: str, target_format: str) -> str:
        input_ext = Path(input_path).suffix.lower().lstrip('.')
        target_format = target_format.lower().lstrip('.')
        conversion_map = {
            ('pdf', 'docx'): self.convert_pdf_to_word,
            ('pdf', 'xlsx'): self.convert_pdf_to_excel,
            ('pdf', 'pptx'): self.convert_pdf_to_powerpoint,
            ('docx', 'pdf'): self.convert_word_to_pdf,
            ('docx', 'xlsx'): self.convert_word_to_excel,
            ('docx', 'csv'): self.convert_word_to_csv,
            ('xlsx', 'pdf'): self.convert_excel_to_pdf,
            ('xlsx', 'csv'): self.convert_excel_to_csv,
            ('xlsx', 'docx'): self.convert_excel_to_word,
            ('pptx', 'pdf'): self.convert_powerpoint_to_pdf,
            ('pptx', 'docx'): self.convert_powerpoint_to_word,
            ('csv', 'xlsx'): self.convert_csv_to_excel,
            ('txt', 'docx'): self.convert_txt_to_word,
            ('txt', 'pdf'): self.convert_txt_to_pdf,
        }
        converter = conversion_map.get((input_ext, target_format))
        if not converter:
            raise ValueError(f"Conversion from {input_ext} to {target_format} not supported")
        return converter(input_path)
    
    def convert_csv_to_excel(self, csv_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(csv_path).stem}.xlsx"
        wb = openpyxl.Workbook()
        ws = wb.active
        with open(csv_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row_idx, row in enumerate(reader, 1):
                for col_idx, value in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
        wb.save(output_path)
        return str(output_path)
    
    def convert_txt_to_word(self, txt_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(txt_path).stem}.docx"
        doc = Document()
        with open(txt_path, 'r', encoding='utf-8') as f:
            content = f.read()
        for para in content.split('\n\n'):
            if para.strip():
                doc.add_paragraph(para)
        doc.save(output_path)
        return str(output_path)
    
    def convert_txt_to_pdf(self, txt_path: str, output_path: str = None) -> str:
        if not output_path:
            output_path = self.temp_dir / f"{Path(txt_path).stem}.pdf"
        with open(txt_path, 'r', encoding='utf-8') as f:
            content = f.read()
        pdf_doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        for para in content.split('\n\n'):
            if para.strip():
                story.append(Paragraph(para.replace('\n', ' '), styles['Normal']))
                story.append(Spacer(1, 12))
        pdf_doc.build(story)
        return str(output_path)
    
    # ============================================================================
    # Document Translation
    # ============================================================================
    
    async def translate_document(self, input_path: str, source_lang: str, target_lang: str, output_format: str = None) -> str:
        text = await self._extract_text(input_path)
        if not text:
            raise ValueError("Could not extract text from document")
        translated_text = await self._translate_text(text, source_lang, target_lang)
        output_format = output_format or Path(input_path).suffix.lower().lstrip('.')
        output_path = self.temp_dir / f"{Path(input_path).stem}_{source_lang}_to_{target_lang}.{output_format}"
        output_path = await self._create_document_from_text(translated_text, output_format, output_path)
        return str(output_path)
    
    async def _extract_text(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == '.docx':
            doc = Document(file_path)
            return '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
        elif ext == '.pdf':
            text = []
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text.append(page.extract_text())
            return '\n'.join(text)
        elif ext in ['.xlsx', '.xls']:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            lines = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                lines.append(f"\n[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    row_text = ' | '.join([str(cell) if cell else '' for cell in row])
                    if row_text.strip():
                        lines.append(row_text)
            return '\n'.join(lines)
        elif ext == '.pptx':
            prs = Presentation(file_path)
            lines = []
            for slide_num, slide in enumerate(prs.slides, 1):
                lines.append(f"\n[Slide {slide_num}]")
                if slide.shapes.title:
                    lines.append(slide.shapes.title.text)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                        lines.append(shape.text)
            return '\n'.join(lines)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    
    async def _translate_text(self, text: str, source_lang: str, target_lang: str) -> str:
        if not self.deepseek_api_key:
            raise ValueError("DeepSeek API key not configured for translation")
        chunks = [text[i:i+4000] for i in range(0, len(text), 4000)]
        translated = []
        for chunk in chunks:
            prompt = f"""Translate the following text from {source_lang} to {target_lang}.
Keep the original formatting, line breaks, and structure.
Only return the translated text, no explanations.

TEXT:
{chunk}

TRANSLATED TEXT:
"""
            url = f"{self.deepseek_base_url}/chat/completions"
            headers = {"Authorization": f"Bearer {self.deepseek_api_key}", "Content-Type": "application/json"}
            data = {
                "model": "deepseek-chat",
                "messages": [
                    {"role": "system", "content": f"You are a professional translator. Translate from {source_lang} to {target_lang} accurately."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.3,
                "max_tokens": min(len(chunk) * 2, 8000)
            }
            async with httpx.AsyncClient(timeout=60.0) as client:
                resp = await client.post(url, headers=headers, json=data)
                if resp.status_code == 200:
                    result = resp.json()
                    translated.append(result["choices"][0]["message"]["content"])
                else:
                    translated.append(chunk)
        return '\n'.join(translated)
    
    async def _create_document_from_text(self, text: str, fmt: str, out_path: Path) -> Path:
        fmt = fmt.lower().lstrip('.')
        if fmt == 'txt':
            with open(out_path, 'w', encoding='utf-8') as f:
                f.write(text)
        elif fmt == 'docx':
            doc = Document()
            for para in text.split('\n\n'):
                if para.strip():
                    doc.add_paragraph(para)
            doc.save(out_path)
        elif fmt == 'pdf':
            pdf_doc = SimpleDocTemplate(str(out_path), pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            for para in text.split('\n\n'):
                if para.strip():
                    story.append(Paragraph(para.replace('\n', ' '), styles['Normal']))
                    story.append(Spacer(1, 12))
            pdf_doc.build(story)
        else:
            # default to word
            out_path = out_path.with_suffix('.docx')
            doc = Document()
            for para in text.split('\n\n'):
                if para.strip():
                    doc.add_paragraph(para)
            doc.save(out_path)
        return out_path
    
    # ============================================================================
    # Batch Conversion
    # ============================================================================
    
    async def batch_convert(self, input_paths: List[str], target_format: str) -> List[str]:
        output_paths = []
        for path in input_paths:
            try:
                out = self.convert_document(path, target_format)
                output_paths.append(out)
                logger.info(f"Converted: {path} -> {out}")
            except Exception as e:
                logger.error(f"Failed to convert {path}: {e}")
        return output_paths


# Global instance
office_editor = OfficeEditor()