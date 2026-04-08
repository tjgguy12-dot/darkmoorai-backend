"""
DarkmoorAI Office Suite - Document Editor
Create, edit, and convert Word, Excel, PowerPoint documents
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Optional
from datetime import datetime
import io
import base64

# Word processing
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

# Excel processing
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference, LineChart, PieChart
from openpyxl.utils import get_column_letter

# PowerPoint processing
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor

# PDF processing
import PyPDF2
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

from app.utils.logger import logger


class OfficeEditor:
    """Office document editor suite"""
    
    def __init__(self):
        self.temp_dir = Path("./data/office_temp")
        self.temp_dir.mkdir(parents=True, exist_ok=True)
    
    # ============================================================================
    # WORD DOCUMENTS
    # ============================================================================
    
    def create_word_document(self, content: Dict[str, Any], filename: str) -> str:
        """Create a Word document"""
        
        doc = Document()
        
        # Add title
        if content.get('title'):
            title = doc.add_heading(content['title'], 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add author and date
        if content.get('author') or content.get('date'):
            info_para = doc.add_paragraph()
            if content.get('author'):
                info_para.add_run(f"Author: {content['author']}").bold = True
                info_para.add_run("\n")
            if content.get('date'):
                info_para.add_run(f"Date: {content.get('date', datetime.now().strftime('%Y-%m-%d'))}")
            info_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            doc.add_paragraph()
        
        # Add sections
        for section in content.get('sections', []):
            # Add heading
            doc.add_heading(section.get('heading', ''), level=1)
            
            # Add content
            for para in section.get('paragraphs', []):
                p = doc.add_paragraph()
                run = p.add_run(para.get('text', ''))
                if para.get('bold'):
                    run.bold = True
                if para.get('italic'):
                    run.italic = True
                if para.get('font_size'):
                    run.font.size = Pt(para['font_size'])
            
            # Add bullet points
            for bullet in section.get('bullets', []):
                p = doc.add_paragraph(bullet, style='List Bullet')
            
            # Add numbered list
            for num, item in enumerate(section.get('numbered_list', []), 1):
                p = doc.add_paragraph(f"{num}. {item}", style='List Number')
            
            # Add table
            if section.get('table'):
                table_data = section['table']
                rows = len(table_data.get('rows', []))
                cols = len(table_data.get('columns', []))
                table = doc.add_table(rows=rows, cols=cols)
                table.style = 'Table Grid'
                
                # Add headers
                for i, col in enumerate(table_data.get('columns', [])):
                    table.cell(0, i).text = col
                    table.cell(0, i).paragraphs[0].runs[0].bold = True
                
                # Add data
                for i, row in enumerate(table_data.get('rows', [])):
                    for j, cell in enumerate(row):
                        table.cell(i + 1, j).text = str(cell)
            
            doc.add_paragraph()  # Add spacing
        
        # Save document
        filepath = self.temp_dir / filename
        doc.save(filepath)
        return str(filepath)
    
    def edit_word_document(self, filepath: str, edits: Dict[str, Any]) -> str:
        """Edit existing Word document"""
        
        doc = Document(filepath)
        
        # Edit paragraphs
        for para_edit in edits.get('paragraphs', []):
            idx = para_edit.get('index')
            text = para_edit.get('text')
            if idx is not None and idx < len(doc.paragraphs):
                doc.paragraphs[idx].text = text
        
        # Add new paragraph
        if edits.get('add_paragraph'):
            doc.add_paragraph(edits['add_paragraph'])
        
        # Save
        output_path = self.temp_dir / f"edited_{Path(filepath).name}"
        doc.save(output_path)
        return str(output_path)
    
    def convert_word_to_pdf(self, word_path: str) -> str:
        """Convert Word to PDF"""
        # Note: This requires additional libraries in production
        # For now, return the original path
        return word_path
    
    # ============================================================================
    # EXCEL / SPREADSHEETS
    # ============================================================================
    
    def create_excel_document(self, content: Dict[str, Any], filename: str) -> str:
        """Create Excel spreadsheet"""
        
        wb = openpyxl.Workbook()
        
        # Create sheets
        for sheet_data in content.get('sheets', []):
            sheet_name = sheet_data.get('name', 'Sheet')
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
            else:
                ws = wb.create_sheet(sheet_name)
            
            # Add headers
            headers = sheet_data.get('headers', [])
            if headers:
                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col, value=header)
                    cell.font = Font(bold=True, size=12)
                    cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
            
            # Add data
            for row_idx, row_data in enumerate(sheet_data.get('data', []), start=2):
                for col_idx, value in enumerate(row_data, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
            
            # Add formulas
            for formula in sheet_data.get('formulas', []):
                cell = formula.get('cell')
                formula_text = formula.get('formula')
                if cell and formula_text:
                    ws[cell] = formula_text
            
            # Add charts
            for chart_data in sheet_data.get('charts', []):
                if chart_data.get('type') == 'bar':
                    self._add_bar_chart(ws, chart_data)
                elif chart_data.get('type') == 'line':
                    self._add_line_chart(ws, chart_data)
            
            # Auto-size columns
            for column in ws.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
        
        # Remove default sheet if we created custom ones
        if len(content.get('sheets', [])) > 0 and 'Sheet' in wb.sheetnames:
            wb.remove(wb['Sheet'])
        
        # Save
        filepath = self.temp_dir / filename
        wb.save(filepath)
        return str(filepath)
    
    def _add_bar_chart(self, ws, chart_data: Dict):
        """Add bar chart to worksheet"""
        
        chart = BarChart()
        chart.title = chart_data.get('title', 'Chart')
        chart.x_axis.title = chart_data.get('x_title', 'Categories')
        chart.y_axis.title = chart_data.get('y_title', 'Values')
        
        data = Reference(ws, min_col=2, min_row=1, max_row=chart_data.get('rows', 5), max_col=chart_data.get('cols', 3))
        categories = Reference(ws, min_col=1, min_row=2, max_row=chart_data.get('rows', 5))
        
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(categories)
        
        ws.add_chart(chart, chart_data.get('position', 'E2'))
    
    def _add_line_chart(self, ws, chart_data: Dict):
        """Add line chart to worksheet"""
        
        chart = LineChart()
        chart.title = chart_data.get('title', 'Line Chart')
        
        data = Reference(ws, min_col=2, min_row=1, max_row=chart_data.get('rows', 5), max_col=chart_data.get('cols', 3))
        categories = Reference(ws, min_col=1, min_row=2, max_row=chart_data.get('rows', 5))
        
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(categories)
        
        ws.add_chart(chart, chart_data.get('position', 'E2'))
    
    def edit_excel_document(self, filepath: str, edits: Dict[str, Any]) -> str:
        """Edit Excel spreadsheet"""
        
        wb = openpyxl.load_workbook(filepath)
        
        for sheet_edit in edits.get('sheets', []):
            sheet_name = sheet_edit.get('name')
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                # Update cells
                for cell_update in sheet_edit.get('cells', []):
                    cell = cell_update.get('cell')
                    value = cell_update.get('value')
                    if cell and value:
                        ws[cell] = value
                
                # Add new rows
                for row_data in sheet_edit.get('add_rows', []):
                    ws.append(row_data)
        
        output_path = self.temp_dir / f"edited_{Path(filepath).name}"
        wb.save(output_path)
        return str(output_path)
    
    def convert_excel_to_csv(self, excel_path: str, sheet_name: str = None) -> str:
        """Convert Excel to CSV"""
        
        import csv
        
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        ws = wb[sheet_name] if sheet_name else wb.active
        
        csv_path = self.temp_dir / f"{Path(excel_path).stem}.csv"
        
        with open(csv_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            for row in ws.iter_rows(values_only=True):
                writer.writerow(row)
        
        return str(csv_path)
    
    # ============================================================================
    # POWERPOINT PRESENTATIONS
    # ============================================================================
    
    def create_powerpoint(self, content: Dict[str, Any], filename: str) -> str:
        """Create PowerPoint presentation"""
        
        prs = Presentation()
        
        # Add title slide
        if content.get('title_slide'):
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = content['title_slide'].get('title', 'Presentation')
            subtitle.text = content['title_slide'].get('subtitle', '')
        
        # Add content slides
        for slide_data in content.get('slides', []):
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = slide_data.get('title', 'Slide')
            
            # Content
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            
            for para in slide_data.get('content', []):
                p = text_frame.add_paragraph()
                p.text = para.get('text', '')
                if para.get('bullet'):
                    p.level = 0
                    p.text = f"• {p.text}"
                if para.get('bold'):
                    p.font.bold = True
            
            # Add image if specified
            if slide_data.get('image'):
                left = PptxInches(1)
                top = PptxInches(2)
                pic = slide.shapes.add_picture(slide_data['image'], left, top, width=PptxInches(6))
        
        # Save
        filepath = self.temp_dir / filename
        prs.save(filepath)
        return str(filepath)
    
    def convert_powerpoint_to_pdf(self, pptx_path: str) -> str:
        """Convert PowerPoint to PDF"""
        # Note: This requires additional libraries
        return pptx_path
    
    # ============================================================================
    # CONVERSION UTILITIES
    # ============================================================================
    
    def convert_document(self, input_path: str, output_format: str) -> str:
        """Convert between document formats"""
        
        input_ext = Path(input_path).suffix.lower()
        output_path = self.temp_dir / f"{Path(input_path).stem}.{output_format}"
        
        # Word to PDF
        if input_ext in ['.docx', '.doc'] and output_format == 'pdf':
            return self.convert_word_to_pdf(input_path)
        
        # Excel to CSV
        elif input_ext in ['.xlsx', '.xls'] and output_format == 'csv':
            return self.convert_excel_to_csv(input_path)
        
        # PDF to Text
        elif input_ext == '.pdf' and output_format == 'txt':
            return self.convert_pdf_to_text(input_path)
        
        else:
            return input_path
    
    def convert_pdf_to_text(self, pdf_path: str) -> str:
        """Extract text from PDF"""
        
        text_content = []
        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text_content.append(page.extract_text())
        
        txt_path = self.temp_dir / f"{Path(pdf_path).stem}.txt"
        with open(txt_path, 'w', encoding='utf-8') as f:
            f.write('\n\n'.join(text_content))
        
        return str(txt_path)
    
    def create_from_template(self, template_type: str, data: Dict) -> str:
        """Create document from template"""
        
        templates = {
            'invoice': self._create_invoice,
            'report': self._create_report,
            'resume': self._create_resume,
            'business_letter': self._create_business_letter,
            'presentation': self._create_presentation_template,
            'budget': self._create_budget_spreadsheet
        }
        
        creator = templates.get(template_type)
        if creator:
            return creator(data)
        return None
    
    def _create_invoice(self, data: Dict) -> str:
        """Create invoice document"""
        
        doc = Document()
        
        # Header
        title = doc.add_heading('INVOICE', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Invoice details
        doc.add_paragraph(f"Invoice Number: {data.get('number', 'INV-001')}")
        doc.add_paragraph(f"Date: {data.get('date', datetime.now().strftime('%Y-%m-%d'))}")
        doc.add_paragraph(f"Due Date: {data.get('due_date', '')}")
        doc.add_paragraph()
        
        # Bill to
        doc.add_heading('Bill To:', level=2)
        doc.add_paragraph(data.get('customer', ''))
        doc.add_paragraph()
        
        # Items table
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Table Grid'
        headers = table.rows[0].cells
        headers[0].text = 'Item'
        headers[1].text = 'Quantity'
        headers[2].text = 'Unit Price'
        headers[3].text = 'Total'
        
        total = 0
        for item in data.get('items', []):
            row = table.add_row().cells
            row[0].text = item.get('description', '')
            row[1].text = str(item.get('quantity', 1))
            row[2].text = f"${item.get('unit_price', 0):.2f}"
            item_total = item.get('quantity', 1) * item.get('unit_price', 0)
            row[3].text = f"${item_total:.2f}"
            total += item_total
        
        # Total
        doc.add_paragraph()
        total_para = doc.add_paragraph()
        total_para.add_run(f"Total: ${total:.2f}").bold = True
        total_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        
        filepath = self.temp_dir / f"invoice_{data.get('number', '001')}.docx"
        doc.save(filepath)
        return str(filepath)
    
    def _create_report(self, data: Dict) -> str:
        """Create report document"""
        
        doc = Document()
        
        # Title
        title = doc.add_heading(data.get('title', 'Report'), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Executive summary
        doc.add_heading('Executive Summary', level=1)
        doc.add_paragraph(data.get('summary', ''))
        
        # Sections
        for section in data.get('sections', []):
            doc.add_heading(section.get('title', ''), level=2)
            doc.add_paragraph(section.get('content', ''))
        
        # Conclusion
        doc.add_heading('Conclusion', level=1)
        doc.add_paragraph(data.get('conclusion', ''))
        
        filepath = self.temp_dir / f"report_{data.get('title', 'report')}.docx"
        doc.save(filepath)
        return str(filepath)
    
    def _create_resume(self, data: Dict) -> str:
        """Create resume document"""
        
        doc = Document()
        
        # Name
        name = doc.add_heading(data.get('name', 'Resume'), 0)
        name.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Contact
        contact = doc.add_paragraph()
        contact.add_run(data.get('email', ''))
        contact.add_run(" | ")
        contact.add_run(data.get('phone', ''))
        contact.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()
        
        # Summary
        doc.add_heading('Professional Summary', level=1)
        doc.add_paragraph(data.get('summary', ''))
        
        # Experience
        doc.add_heading('Work Experience', level=1)
        for exp in data.get('experience', []):
            p = doc.add_paragraph()
            p.add_run(exp.get('title', '')).bold = True
            p.add_run(f" - {exp.get('company', '')}")
            doc.add_paragraph(exp.get('description', ''), style='List Bullet')
        
        # Education
        doc.add_heading('Education', level=1)
        for edu in data.get('education', []):
            p = doc.add_paragraph()
            p.add_run(edu.get('degree', '')).bold = True
            doc.add_paragraph(edu.get('school', ''))
        
        # Skills
        doc.add_heading('Skills', level=1)
        skills = data.get('skills', [])
        if skills:
            doc.add_paragraph(', '.join(skills))
        
        filepath = self.temp_dir / f"resume_{data.get('name', 'resume')}.docx"
        doc.save(filepath)
        return str(filepath)
    
    def _create_business_letter(self, data: Dict) -> str:
        """Create business letter"""
        
        doc = Document()
        
        # Sender
        doc.add_paragraph(data.get('sender', ''))
        doc.add_paragraph()
        
        # Date
        doc.add_paragraph(datetime.now().strftime('%B %d, %Y'))
        doc.add_paragraph()
        
        # Recipient
        doc.add_paragraph(data.get('recipient', ''))
        doc.add_paragraph()
        
        # Subject
        doc.add_paragraph(f"RE: {data.get('subject', '')}")
        doc.add_paragraph()
        
        # Body
        doc.add_paragraph(data.get('body', ''))
        doc.add_paragraph()
        
        # Closing
        doc.add_paragraph(f"Sincerely,\n\n{data.get('sender_name', '')}")
        
        filepath = self.temp_dir / f"letter_{data.get('subject', 'letter')}.docx"
        doc.save(filepath)
        return str(filepath)
    
    def _create_presentation_template(self, data: Dict) -> str:
        """Create presentation template"""
        
        prs = Presentation()
        
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = data.get('title', 'Presentation')
        
        # Content slides
        for i in range(data.get('num_slides', 5)):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Slide {i + 1}"
        
        filepath = self.temp_dir / f"presentation_{data.get('title', 'presentation')}.pptx"
        prs.save(filepath)
        return str(filepath)
    
    def _create_budget_spreadsheet(self, data: Dict) -> str:
        """Create budget spreadsheet"""
        
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Budget"
        
        # Headers
        headers = ['Category', 'Budgeted', 'Actual', 'Variance']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
        
        # Data
        for row, category in enumerate(data.get('categories', []), 2):
            ws.cell(row=row, column=1, value=category.get('name', ''))
            ws.cell(row=row, column=2, value=category.get('budgeted', 0))
            ws.cell(row=row, column=3, value=category.get('actual', 0))
            ws.cell(row=row, column=4, value=f"=C{row}-B{row}")
        
        # Totals
        total_row = len(data.get('categories', [])) + 2
        ws.cell(row=total_row, column=1, value="TOTAL").font = Font(bold=True)
        ws.cell(row=total_row, column=2, value=f"=SUM(B2:B{total_row-1})")
        ws.cell(row=total_row, column=3, value=f"=SUM(C2:C{total_row-1})")
        ws.cell(row=total_row, column=4, value=f"=C{total_row}-B{total_row}")
        
        filepath = self.temp_dir / f"budget_{data.get('name', 'budget')}.xlsx"
        wb.save(filepath)
        return str(filepath)


# Global instance
office_editor = OfficeEditor()